from pptx import Presentation
from pptx.util import Inches, Pt, Emu

prs = Presentation(r"c:\Users\kc10g\OneDrive\Documents\JCD Enterprises\V2X powerpoint template.pptx")

print(f"Slide width: {prs.slide_width}, height: {prs.slide_height}")
print(f"Slide width in inches: {prs.slide_width / 914400:.3f}")
print(f"Slide height in inches: {prs.slide_height / 914400:.3f}")
print(f"Number of slides: {len(prs.slides)}")
print(f"Number of slide layouts: {len(prs.slide_layouts)}")
print()

for i, layout in enumerate(prs.slide_layouts):
    print(f"Layout {i}: '{layout.name}'")
    for ph in layout.placeholders:
        print(f"  Placeholder {ph.placeholder_format.idx}: {ph.name} ({ph.placeholder_format.type}), "
              f"pos=({ph.left},{ph.top}), size=({ph.width},{ph.height})")
    for shape in layout.shapes:
        if not shape.is_placeholder:
            print(f"  Shape: type={shape.shape_type}, name='{shape.name}', "
                  f"pos=({shape.left},{shape.top}), size=({shape.width},{shape.height})")
            try:
                if hasattr(shape, "image"):
                    print(f"    -> Image: {shape.image.content_type}")
            except:
                pass
            if shape.has_text_frame:
                text = shape.text_frame.text[:80]
                if text.strip():
                    print(f"    -> Text: '{text}'")
    print()

print("=== Slide Master ===")
master = prs.slide_masters[0]
for shape in master.shapes:
    desc = f"  Shape: type={shape.shape_type}, name='{shape.name}', pos=({shape.left},{shape.top}), size=({shape.width},{shape.height})"
    print(desc)
    try:
        if hasattr(shape, "image"):
            print(f"    -> Image: {shape.image.content_type}")
    except:
        pass
    if shape.has_text_frame:
        text = shape.text_frame.text[:100]
        if text.strip():
            print(f"    -> Text: '{text}'")

print()
print("=== Actual Slides ===")
for si, slide in enumerate(prs.slides):
    print(f"Slide {si+1}: layout='{slide.slide_layout.name}'")
    for shape in slide.shapes:
        desc = f"  Shape: type={shape.shape_type}, name='{shape.name}', pos=({shape.left},{shape.top}), size=({shape.width},{shape.height})"
        print(desc)
        try:
            if hasattr(shape, "image"):
                print(f"    -> Image: {shape.image.content_type}")
        except:
            pass
        if shape.has_text_frame:
            text = shape.text_frame.text[:100]
            if text.strip():
                print(f"    -> Text: '{text}'")
