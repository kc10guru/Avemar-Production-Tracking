import xlrd
from copy import deepcopy
from datetime import datetime, timedelta
from collections import defaultdict
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData

TEMPLATE = (
    r"c:\Users\kc10g\OneDrive\Documents\JCD Enterprises"
    r"\V2X powerpoint template.pptx"
)
SPREADSHEET = (
    r"c:\Users\kc10g\OneDrive\Documents\JCD Enterprises"
    r"\INTERNAL REPORT -VERTEX WINDOWS 2-13-26.xls"
)
OUTPUT = (
    r"c:\Users\kc10g\OneDrive\Documents\JCD Enterprises"
    r"\Vertex Windows - Expansion Proposal.pptx"
)

DARK_NAVY = RGBColor(0x1B, 0x2A, 0x4A)
ACCENT_BLUE = RGBColor(0x00, 0x4B, 0x87)
ACCENT_GREEN = RGBColor(0x27, 0xAE, 0x60)
ACCENT_ORANGE = RGBColor(0xE6, 0x7E, 0x22)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF2, 0xF4, 0xF5)
MED_GRAY = RGBColor(0x85, 0x92, 0x9E)
TEXT_DARK = RGBColor(0x2C, 0x3E, 0x50)

TITLE_LAYOUT_IDX = 0
CONTENT_LAYOUT_IDX = 5
CONTENT_LARGE_TITLE_IDX = 6


def extract_monthly_deliveries():
    wb = xlrd.open_workbook(SPREADSHEET)
    ws = wb.sheet_by_name("QUOTES APPROVED")
    monthly = defaultdict(int)
    for r in range(19, ws.nrows):
        status = str(ws.cell_value(r, 8)).strip().upper()
        shipped_val = ws.cell_value(r, 10)
        if "SHIPPED" in status and isinstance(shipped_val, float) and shipped_val > 0:
            dt = datetime(1899, 12, 30) + timedelta(days=int(shipped_val))
            if dt > datetime(2026, 2, 24):
                dt = dt.replace(year=2025)
            key = dt.strftime("%b %Y")
            monthly[key] += 1
    ordered_keys = sorted(monthly.keys(), key=lambda k: datetime.strptime(k, "%b %Y"))
    return ordered_keys, [monthly[k] for k in ordered_keys]


# ── Helpers ─────────────────────────────────────────────────────────────
def add_text_box(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=TEXT_DARK, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def add_bullet_list(slide, left, top, width, height, items,
                    font_size=16, color=TEXT_DARK, spacing=Pt(8)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = spacing
        if isinstance(item, tuple):
            label, value = item
            r1 = p.add_run()
            r1.text = label
            r1.font.size = Pt(font_size)
            r1.font.bold = True
            r1.font.color.rgb = color
            r2 = p.add_run()
            r2.text = value
            r2.font.size = Pt(font_size)
            r2.font.color.rgb = color
        else:
            r = p.add_run()
            r.text = f"  {item}"
            r.font.size = Pt(font_size)
            r.font.color.rgb = color
    return txBox


def add_kpi_card(slide, left, top, width, height, label, value, accent=ACCENT_BLUE):
    card = slide.shapes.add_shape(1, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = RGBColor(0xD5, 0xDB, 0xDB)
    card.line.width = Pt(1)
    card.shadow.inherit = False
    tf = card.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    p_val = tf.paragraphs[0]
    p_val.alignment = PP_ALIGN.CENTER
    r = p_val.add_run()
    r.text = value
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.color.rgb = accent
    p_label = tf.add_paragraph()
    p_label.alignment = PP_ALIGN.CENTER
    r2 = p_label.add_run()
    r2.text = label
    r2.font.size = Pt(12)
    r2.font.color.rgb = MED_GRAY


def set_slide_title(slide, text):
    """Set the title placeholder text if it exists."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = ""
            p = ph.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = text
            run.font.size = Pt(24)
            run.font.bold = True
            run.font.color.rgb = DARK_NAVY
            return
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 10:
            ph.text = ""
            p = ph.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = text
            run.font.size = Pt(24)
            run.font.bold = True
            run.font.color.rgb = DARK_NAVY
            return


def clear_body_placeholder(slide):
    """Clear body placeholders so we can add our own content."""
    for ph in slide.placeholders:
        idx = ph.placeholder_format.idx
        if idx in (1, 10, 11, 12, 13, 14):
            sp = ph._element
            sp.getparent().remove(sp)


def delete_all_slides(prs):
    """Remove all existing slides from the template."""
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].get(
            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
        )
        prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])


# ── Build ───────────────────────────────────────────────────────────────
def build():
    prs = Presentation(TEMPLATE)
    delete_all_slides(prs)

    title_layout = prs.slide_layouts[TITLE_LAYOUT_IDX]
    content_layout = prs.slide_layouts[CONTENT_LAYOUT_IDX]
    content_lg_layout = prs.slide_layouts[CONTENT_LARGE_TITLE_IDX]

    months, counts = extract_monthly_deliveries()

    # Content area boundaries (from template layout 5)
    CX = Inches(0.56)   # content left
    CY = Inches(1.15)   # content top (below title + line)
    CW = Inches(11.75)  # content width
    CH = Inches(5.8)    # content height

    # ── SLIDE 1: Title ──────────────────────────────────────────────────
    s1 = prs.slides.add_slide(title_layout)
    for ph in s1.placeholders:
        idx = ph.placeholder_format.idx
        if idx == 10:
            ph.text = ""
            p = ph.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            r = p.add_run()
            r.text = "Vertex Windows\nWindshield Repair Program"
            r.font.size = Pt(36)
            r.font.bold = True
            r.font.color.rgb = WHITE
            p2 = ph.text_frame.add_paragraph()
            p2.alignment = PP_ALIGN.LEFT
            r2 = p2.add_run()
            r2.text = "Expansion Proposal"
            r2.font.size = Pt(28)
            r2.font.bold = False
            r2.font.color.rgb = RGBColor(0x7F, 0xB3, 0xD8)
        elif idx == 11:
            ph.text = ""
            p = ph.text_frame.paragraphs[0]
            r = p.add_run()
            r.text = "Prepared for V2X Leadership  |  February 2026"
            r.font.size = Pt(14)
            r.font.color.rgb = RGBColor(0xBD, 0xBD, 0xBD)

    # ── SLIDE 2: Past Performance ───────────────────────────────────────
    s2 = prs.slides.add_slide(content_layout)
    set_slide_title(s2, "Past Performance \u2014 Units Delivered Per Month")
    clear_body_placeholder(s2)

    chart_data = CategoryChartData()
    chart_data.categories = months
    chart_data.add_series("Units Shipped", counts)

    chart_frame = s2.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        CX, Inches(1.2), Inches(9.2), Inches(4.8), chart_data,
    )
    chart = chart_frame.chart
    chart.has_legend = False
    chart.style = 2
    plot = chart.plots[0]
    plot.gap_width = 80
    series = plot.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = ACCENT_BLUE
    series.has_data_labels = True
    series.data_labels.font.size = Pt(12)
    series.data_labels.font.bold = True
    series.data_labels.font.color.rgb = DARK_NAVY
    series.data_labels.number_format = "0"
    series.data_labels.label_position = XL_LABEL_POSITION.OUTSIDE_END
    chart.category_axis.tick_labels.font.size = Pt(11)
    chart.category_axis.tick_labels.font.color.rgb = TEXT_DARK
    chart.value_axis.visible = False
    chart.value_axis.has_major_gridlines = False

    # Sidebar with monthly totals
    sidebar_left = Inches(9.9)
    sidebar_top = Inches(1.2)
    sidebar_w = Inches(2.8)
    total_shipped = sum(counts)
    avg = total_shipped / len(counts) if counts else 0

    sidebar_box = s2.shapes.add_shape(
        1, sidebar_left, sidebar_top, sidebar_w, Inches(4.8)
    )
    sidebar_box.fill.solid()
    sidebar_box.fill.fore_color.rgb = LIGHT_GRAY
    sidebar_box.line.color.rgb = ACCENT_BLUE
    sidebar_box.line.width = Pt(1.5)

    add_text_box(s2, sidebar_left + Inches(0.15), sidebar_top + Inches(0.1),
                 sidebar_w - Inches(0.3), Inches(0.35),
                 "Monthly Totals", 15, True, ACCENT_BLUE, PP_ALIGN.CENTER)

    sidebar_tb = s2.shapes.add_textbox(
        sidebar_left + Inches(0.15), sidebar_top + Inches(0.5),
        sidebar_w - Inches(0.3), Inches(3.3),
    )
    tf = sidebar_tb.text_frame
    tf.word_wrap = True
    for i, (m, c) in enumerate(zip(months, counts)):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(3)
        r_m = p.add_run()
        r_m.text = f"{m}:  "
        r_m.font.size = Pt(13)
        r_m.font.color.rgb = TEXT_DARK
        r_c = p.add_run()
        r_c.text = str(c)
        r_c.font.size = Pt(13)
        r_c.font.bold = True
        r_c.font.color.rgb = ACCENT_BLUE

    # Divider line + totals at bottom of sidebar
    p_div = tf.add_paragraph()
    p_div.space_before = Pt(6)
    p_div.space_after = Pt(4)
    r_div = p_div.add_run()
    r_div.text = "\u2500" * 20
    r_div.font.size = Pt(8)
    r_div.font.color.rgb = MED_GRAY

    p_total = tf.add_paragraph()
    r_tl = p_total.add_run()
    r_tl.text = "Total:  "
    r_tl.font.size = Pt(14)
    r_tl.font.bold = True
    r_tl.font.color.rgb = TEXT_DARK
    r_tv = p_total.add_run()
    r_tv.text = str(total_shipped)
    r_tv.font.size = Pt(14)
    r_tv.font.bold = True
    r_tv.font.color.rgb = ACCENT_GREEN

    p_avg = tf.add_paragraph()
    r_al = p_avg.add_run()
    r_al.text = "Avg/Mo:  "
    r_al.font.size = Pt(13)
    r_al.font.color.rgb = TEXT_DARK
    r_av = p_avg.add_run()
    r_av.text = f"{avg:.1f}"
    r_av.font.size = Pt(13)
    r_av.font.bold = True
    r_av.font.color.rgb = ACCENT_BLUE

    # ── SLIDE 3: Current Business Model ─────────────────────────────────
    s3 = prs.slides.add_slide(content_layout)
    set_slide_title(s3, "Current Business Model")
    clear_body_placeholder(s3)

    items_current = [
        ("Overhauled Units:  ", "$40,000 sale  /  $29,500 cost  =  $10,500 profit"),
        ("Repaired Units:  ", "$30,000 sale  /  $19,500 cost  =  $10,500 profit"),
        ("Projected Annual Sales:  ", "25 units"),
        ("Sales Mix:  ", "~90% overhaul  /  ~10% repair"),
    ]
    add_bullet_list(s3, CX, Inches(1.2), Inches(11), Inches(2.8), items_current, 19)

    card_y = Inches(4.5)
    card_w = Inches(3.4)
    card_h = Inches(1.5)
    gap = Inches(0.6)
    add_kpi_card(s3, CX + Inches(0.3), card_y, card_w, card_h,
                 "Annual Revenue", "$975,000", ACCENT_BLUE)
    add_kpi_card(s3, CX + Inches(0.3) + card_w + gap, card_y, card_w, card_h,
                 "Annual Cost", "$712,500", ACCENT_ORANGE)
    add_kpi_card(s3, CX + Inches(0.3) + 2 * (card_w + gap), card_y, card_w, card_h,
                 "Annual Profit", "$262,500", ACCENT_GREEN)

    # ── SLIDE 4: Proposed Expansion Model ───────────────────────────────
    s4 = prs.slides.add_slide(content_layout)
    set_slide_title(s4, "Proposed Expansion Model")
    clear_body_placeholder(s4)

    items_proposed = [
        ("Flat Rate Cost (to shop):  ", "$19,000 per unit"),
        ("Sales Price:  ", "$24,000 per unit"),
        ("Exchange (with core return):  ", "$30,000"),
        ("Outright Sale (no core):  ", "$33,000"),
        ("Annual Target:  ", "100 units / year  (2 per week)"),
        ("Profit Per Unit:  ", "$5,000"),
    ]
    add_bullet_list(s4, CX, Inches(1.2), Inches(11), Inches(3.2), items_proposed, 19)

    add_kpi_card(s4, CX + Inches(0.3), card_y, card_w, card_h,
                 "Annual Revenue", "$2,400,000", ACCENT_BLUE)
    add_kpi_card(s4, CX + Inches(0.3) + card_w + gap, card_y, card_w, card_h,
                 "Annual Cost", "$1,900,000", ACCENT_ORANGE)
    add_kpi_card(s4, CX + Inches(0.3) + 2 * (card_w + gap), card_y, card_w, card_h,
                 "Annual Profit", "$500,000", ACCENT_GREEN)

    # ── SLIDE 5: Production & Financial Structure ───────────────────────
    s5 = prs.slides.add_slide(content_layout)
    set_slide_title(s5, "Production & Financial Structure")
    clear_body_placeholder(s5)

    add_text_box(s5, CX, Inches(1.2), Inches(5.5), Inches(0.4),
                 "Financial Terms", 20, True, ACCENT_BLUE)
    left_items = [
        ("Monthly Material Advance:  ", "$108,000"),
        ("Material Coverage:  ", "10 windshields per month"),
        ("Annual PO to Shop:  ", "$1,080,000  (12 monthly payments)"),
        ("Individual Repair POs:  ", "$8,200 each"),
    ]
    add_bullet_list(s5, CX, Inches(1.8), Inches(5.5), Inches(3.8), left_items, 17)

    add_text_box(s5, Inches(6.8), Inches(1.2), Inches(5.5), Inches(0.4),
                 "Production Terms", 20, True, ACCENT_BLUE)
    right_items = [
        ("Production Rate:  ", "2 units per week"),
        ("Late Delivery Penalty:  ", "$500 per unit"),
        ("Core Return:  ", "Repairable core turn-in required for exchange pricing"),
    ]
    add_bullet_list(s5, Inches(6.8), Inches(1.8), Inches(5.5), Inches(3.8),
                    right_items, 17)

    # ── SLIDE 6: Lead Time & Delivery Schedule ──────────────────────────
    s6 = prs.slides.add_slide(content_layout)
    set_slide_title(s6, "Lead Time & Delivery Schedule")
    clear_body_placeholder(s6)

    phase1_box = s6.shapes.add_shape(
        1, CX, Inches(1.5), Inches(5.4), Inches(3.5)
    )
    phase1_box.fill.solid()
    phase1_box.fill.fore_color.rgb = LIGHT_GRAY
    phase1_box.line.color.rgb = ACCENT_BLUE
    phase1_box.line.width = Pt(2)

    add_text_box(s6, CX + Inches(0.3), Inches(1.7), Inches(4.8), Inches(0.4),
                 "Phase 1 \u2014 First 90 Days", 20, True, ACCENT_BLUE)
    add_bullet_list(s6, CX + Inches(0.3), Inches(2.3), Inches(4.8), Inches(2.5), [
        "Lead time: 60 days (receipt to delivery)",
        "Ramp-up to full production rate",
        "Establish material pipeline",
        "Build initial inventory buffer",
    ], 16)

    phase2_box = s6.shapes.add_shape(
        1, Inches(6.8), Inches(1.5), Inches(5.4), Inches(3.5)
    )
    phase2_box.fill.solid()
    phase2_box.fill.fore_color.rgb = LIGHT_GRAY
    phase2_box.line.color.rgb = ACCENT_GREEN
    phase2_box.line.width = Pt(2)

    add_text_box(s6, Inches(7.1), Inches(1.7), Inches(4.8), Inches(0.4),
                 "Phase 2 \u2014 Steady State", 20, True, ACCENT_GREEN)
    add_bullet_list(s6, Inches(7.1), Inches(2.3), Inches(4.8), Inches(2.5), [
        "Lead time reduced to 45 days",
        "Sustained 2 units per week",
        "Consistent monthly deliveries",
        "Penalty enforcement active",
    ], 16)

    arrow = s6.shapes.add_shape(13, Inches(6.2), Inches(2.9), Inches(0.5), Inches(0.5))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = ACCENT_BLUE
    arrow.line.fill.background()

    # ── SLIDE 7: Annual Financial Projections ───────────────────────────
    s7 = prs.slides.add_slide(content_layout)
    set_slide_title(s7, "Annual Financial Projections \u2014 Current vs. Proposed")
    clear_body_placeholder(s7)

    col_headers = ["", "Current Model", "Proposed Model"]
    row_data = [
        ["Units / Year", "25", "100"],
        ["Sales Mix", "90% OH / 10% Repair", "Flat Rate"],
        ["Sale Price", "$40K OH / $30K Repair", "$24,000"],
        ["Cost Per Unit", "$29.5K OH / $19.5K Repair", "$19,000"],
        ["Profit Per Unit", "$10,500", "$5,000"],
        ["Annual Revenue", "$975,000", "$2,400,000"],
        ["Annual Cost", "$712,500", "$1,900,000"],
        ["Annual Profit", "$262,500", "$500,000"],
    ]

    tbl_left = Inches(1.2)
    tbl_top = Inches(1.3)
    tbl_width = Inches(10.5)
    tbl_height = Inches(4.2)
    rows = len(row_data) + 1
    cols = 3

    table_shape = s7.shapes.add_table(rows, cols, tbl_left, tbl_top, tbl_width, tbl_height)
    table = table_shape.table
    for i, w in enumerate([Inches(3.0), Inches(3.75), Inches(3.75)]):
        table.columns[i].width = w

    for ci, header in enumerate(col_headers):
        cell = table.cell(0, ci)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_NAVY
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                r.font.size = Pt(15)
                r.font.bold = True
                r.font.color.rgb = WHITE

    for ri, row in enumerate(row_data):
        is_profit_row = ri == len(row_data) - 1
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.text = val
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if is_profit_row:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xD5, 0xF5, 0xE3)
            elif ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
                for r in p.runs:
                    r.font.size = Pt(14)
                    r.font.color.rgb = TEXT_DARK
                    if ci == 0:
                        r.font.bold = True
                    if is_profit_row:
                        r.font.bold = True
                        r.font.color.rgb = ACCENT_GREEN
                        r.font.size = Pt(16)

    add_text_box(
        s7, Inches(1.2), Inches(5.8), Inches(10.5), Inches(0.5),
        "Proposed model delivers nearly 2x annual profit "
        "($500K vs $262.5K) with 2.5x revenue growth.",
        15, True, ACCENT_GREEN, PP_ALIGN.CENTER,
    )

    # ── SLIDE 8: Summary & Recommendation ──────────────────────────────
    s8 = prs.slides.add_slide(title_layout)
    for ph in s8.placeholders:
        idx = ph.placeholder_format.idx
        if idx == 10:
            ph.text = ""
            p = ph.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            r = p.add_run()
            r.text = "Summary & Recommendation"
            r.font.size = Pt(32)
            r.font.bold = True
            r.font.color.rgb = WHITE
        elif idx == 11:
            ph.text = ""

    summary_items = [
        "Transition from variable pricing to a flat-rate model at $24K per unit",
        "Scale production from ~3 units/month to 8+ units/month (2 per week)",
        "Increase annual volume from 25 to 100 units",
        "Grow annual profit from $262,500 to $500,000 \u2014 a 90% increase",
        "Structured PO of $1,080,000 provides shop with predictable cash flow",
        "Built-in $500 late-delivery penalty drives on-time performance",
        "Lead times improve from 60 days to 45 days within first 90 days",
    ]
    add_bullet_list(
        s8, Inches(1.2), Inches(3.2), Inches(9.5), Inches(3.0),
        summary_items, 17, WHITE, Pt(10),
    )

    add_text_box(
        s8, Inches(1.2), Inches(5.8), Inches(9.5), Inches(0.7),
        "Recommendation: Approve the expanded flat-rate production model "
        "to capture significantly higher margins at scale.",
        18, True, RGBColor(0x7F, 0xB3, 0xD8), PP_ALIGN.LEFT,
    )

    # ── Save ────────────────────────────────────────────────────────────
    prs.save(OUTPUT)
    print(f"Presentation saved to:\n  {OUTPUT}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    build()
